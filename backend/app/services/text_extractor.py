"""Extract text from common office file formats.

Supports: PDF, DOCX, XLSX, PPTX, CSV
Saves extracted text as a companion .txt file alongside the original.
"""

import io
from pathlib import Path

import httpx
from loguru import logger


# File extensions that need text extraction
EXTRACTABLE_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Text extensions that don't need extraction
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
             ".js", ".ts", ".py", ".html", ".css", ".sh", ".log", ".env"}


def needs_extraction(filename: str) -> bool:
    """Check if a file needs text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in EXTRACTABLE_EXTS


def extract_text(file_bytes: bytes, filename: str) -> str | None:
    """Extract text from a binary file.
    
    Returns extracted text string, or None if extraction fails.
    """
    ext = Path(filename).suffix.lower()
    
    try:
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        elif ext == ".docx":
            return _extract_docx(file_bytes)
        elif ext == ".xlsx":
            return _extract_xlsx(file_bytes)
        elif ext == ".pptx":
            return _extract_pptx(file_bytes)
    except Exception as e:
        logger.error(f"[TextExtractor] Failed to extract from {filename}: {e}")
        return None
    
    return None


def save_extracted_text(save_path: Path, file_bytes: bytes, filename: str) -> Path | None:
    """Extract text and save as a companion .txt file.
    
    For example: report.pdf → report.txt
    Returns the path to the text file, or None if extraction failed.
    """
    text = extract_text(file_bytes, filename)
    if not text or not text.strip():
        return None
    
    txt_path = save_path.parent / f"{save_path.stem}.txt"
    txt_path.write_text(text, encoding="utf-8")
    logger.info(f"[TextExtractor] Extracted {len(text)} chars from {filename} → {txt_path.name}")
    return txt_path


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber
    
    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- 第{i+1}页 ---\n{text.strip()}")
            
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                if table:
                    rows = []
                    for row in table:
                        cells = [str(c or "").strip() for c in row]
                        rows.append(" | ".join(cells))
                    if rows:
                        pages.append("表格:\n" + "\n".join(rows))
    
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    
    doc = Document(io.BytesIO(data))
    parts = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve heading hierarchy
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
    
    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n表格:\n" + "\n".join(rows))
    
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        
        if rows:
            parts.append(f"## 工作表: {sheet}\n" + "\n".join(rows))
    
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    
    prs = Presentation(io.BytesIO(data))
    parts = []
    
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        
        if texts:
            parts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(texts))

    return "\n\n".join(parts)


# ─── OCR Functions ───────────────────────────────────────────────────────────────


async def extract_text_with_ocr(
    file_bytes: bytes,
    filename: str,
    ocr_url: str = "http://localhost:6008/ocr/file",
    ocr_api_key: str = ""
) -> str | None:
    """Extract text from file using OCR service.

    Args:
        file_bytes: File content as bytes
        filename: Original filename
        ocr_url: OCR service URL
        ocr_api_key: OCR service API key

    Returns:
        Extracted text or None if OCR fails
    """
    try:
        files = {"file": (filename, file_bytes)}
        headers = {}
        if ocr_api_key:
            headers["Authorization"] = f"Bearer {ocr_api_key}"

        timeout = 60  # OCR can be slow
        async with httpx.AsyncClient(timeout=timeout) as client:
            resp = await client.post(ocr_url, files=files, headers=headers)

            if resp.status_code != 200:
                logger.error(f"[OCR] Service returned status {resp.status_code}: {resp.text[:200]}")
                return None

            result = resp.json()
            if result.get("success"):
                content = result.get("content", "")
                # content may be a list of page dicts or a plain string
                if isinstance(content, list):
                    parts = []
                    for page in content:
                        if isinstance(page, dict):
                            parts.append(page.get("content", ""))
                        else:
                            parts.append(str(page))
                    content = "\n\n".join(parts)
                logger.info(f"[OCR] Successfully extracted {len(content)} chars from {filename} (pages: {result.get('pages', 0)})")
                return content
            else:
                logger.error(f"[OCR] Service returned failure for {filename}")
                return None

    except Exception as e:
        logger.error(f"[OCR] Failed to process {filename}: {e}")
        return None


def is_pdf_scanned(pdf_bytes: bytes) -> bool:
    """Detect if PDF is scanned by checking for text content.

    A PDF is considered scanned if:
    1. Extractable text is very short (< 100 chars) OR
    2. Text extraction fails completely

    Args:
        pdf_bytes: PDF file content as bytes

    Returns:
        True if PDF appears to be scanned, False otherwise
    """
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_text = ""
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    total_text += text

            # If we can extract very little text, it's likely a scan
            return len(total_text.strip()) < 100
    except Exception:
        # If extraction fails, assume it needs OCR
        return True


async def extract_text_with_ocr_if_needed(
    file_bytes: bytes,
    filename: str,
    use_ocr: bool = False,
    ocr_config: dict | None = None
) -> str | None:
    """Extract text with optional OCR.

    Strategy:
    - If use_ocr=True or ocr_config.ocr_enabled=True, use OCR
    - For PDF: check if scanned first, use OCR if needed
    - For images: always use OCR when enabled
    - For other formats: use normal extraction

    Args:
        file_bytes: File content as bytes
        filename: Original filename
        use_ocr: Force OCR regardless of file type
        ocr_config: OCR config dict with keys: enabled, url, api_key

    Returns:
        Extracted text or None if all methods fail
    """
    ext = Path(filename).suffix.lower()

    # Check OCR configuration
    ocr_enabled = ocr_config.get("ocr_enabled", False) if ocr_config else False
    ocr_url = ocr_config.get("ocr_url", "http://localhost:6008/ocr/file") if ocr_config else ""
    ocr_api_key = ocr_config.get("ocr_api_key", "") if ocr_config else ""

    # Images: always use OCR when enabled
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}:
        if use_ocr or ocr_enabled:
            return await extract_text_with_ocr(file_bytes, filename, ocr_url, ocr_api_key)

    # PDF: check if scanned
    if ext == ".pdf":
        if use_ocr or ocr_enabled:
            # Check if PDF is scanned
            if is_pdf_scanned(file_bytes):
                logger.info(f"[OCR] PDF {filename} appears to be scanned, using OCR")
                return await extract_text_with_ocr(file_bytes, filename, ocr_url, ocr_api_key)
            else:
                logger.info(f"[OCR] PDF {filename} has extractable text, using pdfplumber")
                return extract_text(file_bytes, filename)
        else:
            # OCR not enabled, use normal extraction
            return extract_text(file_bytes, filename)

    # Other formats: use normal extraction
    return extract_text(file_bytes, filename)

